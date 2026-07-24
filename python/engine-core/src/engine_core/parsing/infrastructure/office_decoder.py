"""Bounded PDF and OOXML text extraction adapters."""

from __future__ import annotations

from hashlib import sha256
from io import BytesIO
from zipfile import BadZipFile, ZipFile

from engine_core.parsing.domain.models import NormalizedDocument, PageSpan, SourceType
from engine_core.parsing.shared.errors import DocumentParsingError

MAX_ARCHIVE_ENTRIES = 10_000
MAX_EXPANDED_BYTES = 64 * 1024 * 1024


def decode_pdf(content: bytes) -> NormalizedDocument:
    try:
        import fitz  # type: ignore[import-untyped]
    except ImportError as error:
        raise DocumentParsingError("DOC-E-002", "PDF parser is not installed", 503) from error
    try:
        document = fitz.open(stream=content, filetype="pdf")
        if document.needs_pass:
            raise DocumentParsingError("DOC-V-008", "Encrypted PDF is not supported", 400)
        pages = [page.get_text("text", sort=True) for page in document]
    except DocumentParsingError:
        raise
    except Exception as error:
        raise DocumentParsingError("DOC-V-008", "Invalid PDF document", 400) from error
    return _assemble(content, SourceType.PDF, pages)


def decode_ooxml(content: bytes, source_type: SourceType) -> NormalizedDocument:
    _validate_archive(content)
    try:
        if source_type is SourceType.DOCX:
            from docx import Document

            document = Document(BytesIO(content))
            sections = ["\n".join(p.text for p in document.paragraphs)]
            for table in document.tables:
                sections.append("\n".join("\t".join(cell.text for cell in row.cells) for row in table.rows))
        elif source_type is SourceType.PPTX:
            from pptx import Presentation

            presentation = Presentation(BytesIO(content))
            sections = [
                "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
                for slide in presentation.slides
            ]
        elif source_type is SourceType.XLSX:
            from openpyxl import load_workbook  # type: ignore[import-untyped]

            workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
            sections = [
                "\n".join("\t".join("" if value is None else str(value) for value in row) for row in sheet.values)
                for sheet in workbook.worksheets
            ]
        else:
            raise ValueError("Unsupported OOXML source type")
    except ImportError as error:
        raise DocumentParsingError("DOC-E-002", "Office parser is not installed", 503) from error
    except (BadZipFile, KeyError, ValueError, OSError) as error:
        raise DocumentParsingError("DOC-V-008", "Invalid Office document", 400) from error
    return _assemble(content, source_type, sections)


def _validate_archive(content: bytes) -> None:
    try:
        with ZipFile(BytesIO(content)) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_ENTRIES or sum(item.file_size for item in entries) > MAX_EXPANDED_BYTES:
                raise DocumentParsingError("DOC-V-009", "Office archive exceeds extraction limits", 413)
            if any(item.flag_bits & 1 for item in entries):
                raise DocumentParsingError("DOC-V-008", "Encrypted Office document is not supported", 400)
    except BadZipFile as error:
        raise DocumentParsingError("DOC-V-008", "Invalid Office document", 400) from error


def _assemble(content: bytes, source_type: SourceType, sections: list[str]) -> NormalizedDocument:
    from engine_core.parsing.infrastructure.input_decoder import normalize_text

    parts: list[str] = []
    spans: list[PageSpan] = []
    offset = 0
    for page, raw in enumerate(sections, 1):
        text = normalize_text(raw, allow_empty=True)
        if not text:
            continue
        if parts:
            parts.append("\n\n")
            offset += 2
        parts.append(text)
        spans.append(PageSpan(offset, offset + len(text), page))
        offset += len(text)
    normalized = "".join(parts)
    if not normalized:
        raise DocumentParsingError("DOC-V-001", "Document text must not be blank", 400)
    return NormalizedDocument(source_type, sha256(content).hexdigest(), normalized, tuple(spans))

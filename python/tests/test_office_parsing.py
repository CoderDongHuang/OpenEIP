from collections.abc import Callable
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from engine_core.parsing.domain.models import SourceType
from engine_core.parsing.infrastructure import office_decoder
from engine_core.parsing.infrastructure.input_decoder import DocumentInputDecoder
from engine_core.parsing.shared.errors import DocumentParsingError


def _docx() -> bytes:
    target = BytesIO()
    document = Document()
    document.add_paragraph("Word source alpha")
    document.add_table(rows=1, cols=2).rows[0].cells[0].text = "table cell"
    document.save(target)
    return target.getvalue()


def _pptx() -> bytes:
    target = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide source beta"
    presentation.save(target)
    return target.getvalue()


def _xlsx() -> bytes:
    target = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["invoice", 42])
    workbook.save(target)
    return target.getvalue()


def _pdf() -> bytes:
    document = fitz.open()
    document.new_page().insert_text((72, 72), "PDF page one")
    document.new_page().insert_text((72, 72), "PDF page two")
    return document.tobytes()


@pytest.mark.parametrize(
    ("content_type", "factory", "source_type", "expected"),
    [
        ("application/pdf", _pdf, SourceType.PDF, "PDF page two"),
        (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            _docx,
            SourceType.DOCX,
            "Word source alpha",
        ),
        (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            _pptx,
            SourceType.PPTX,
            "Slide source beta",
        ),
        (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            _xlsx,
            SourceType.XLSX,
            "invoice\t42",
        ),
    ],
)
def test_real_document_adapters_extract_traceable_text(
    content_type: str, factory: Callable[[], bytes], source_type: SourceType, expected: str
) -> None:
    content = factory()
    result = DocumentInputDecoder(2 * 1024 * 1024).decode(content, content_type)

    assert result.source_type is source_type
    assert expected in result.text
    assert result.page_spans
    assert all(span.page >= 1 and span.end > span.start for span in result.page_spans)


def test_corrupt_and_expansion_bomb_office_inputs_fail_closed(monkeypatch: pytest.MonkeyPatch) -> None:
    decoder = DocumentInputDecoder(4096)
    media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    with pytest.raises(DocumentParsingError, match="Invalid Office") as corrupt:
        decoder.decode(b"not-a-zip", media_type)
    assert corrupt.value.code == "DOC-V-008"

    target = BytesIO()
    with ZipFile(target, "w", ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", "x" * 32)
    monkeypatch.setattr(office_decoder, "MAX_EXPANDED_BYTES", 16)
    with pytest.raises(DocumentParsingError, match="exceeds extraction limits") as oversized:
        decoder.decode(target.getvalue(), media_type)
    assert oversized.value.code == "DOC-V-009"
